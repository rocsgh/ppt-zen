#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""Assemble a folder of page images into an image-based .pptx (each image full-bleed 16:9).\n\nNon-16:9 images (most models emit 3:2, e.g. 1536x1024) are CENTER COVER-CROPPED to 16:9\ninside the pptx - so keep titles and key content clear of the top/bottom ~8% of the frame.\nNote the output is image-based: text is not editable afterwards; to fix a typo, regenerate\nthat one page.

This is the "how do I turn N generated slides into a deck" step. PPT-Zen decides
and generates the pages; this makes the deliverable.

Usage:
  pip install python-pptx
  python3 scripts/assemble_pptx.py slides/ deck.pptx
"""
import os, sys, glob


def main():
    if len(sys.argv) < 3:
        sys.exit("usage: assemble_pptx.py <images_dir> <out.pptx>")
    src, out = sys.argv[1], sys.argv[2]
    try:
        from pptx import Presentation
        from pptx.util import Inches
    except ImportError:
        sys.exit("Need python-pptx:  pip install python-pptx")
    imgs = sorted(g for ext in ("*.jpg", "*.jpeg", "*.png") for g in glob.glob(os.path.join(src, ext)))
    if not imgs:
        sys.exit("no .jpg/.png images found in " + src)
    prs = Presentation()
    prs.slide_width = Inches(13.333)   # 16:9
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    try:
        from PIL import Image  # python-pptx already depends on Pillow
    except ImportError:
        Image = None
    target = 16.0 / 9.0
    for p in imgs:
        s = prs.slides.add_slide(blank)
        pic = s.shapes.add_picture(p, 0, 0, width=prs.slide_width, height=prs.slide_height)
        if Image is not None:
            with Image.open(p) as im:
                ar = im.width / im.height
            if abs(ar - target) > 0.01:
                if ar < target:   # image too tall (e.g. 3:2) -> crop top/bottom equally
                    keep = ar / target
                    c = (1.0 - keep) / 2.0
                    pic.crop_top = c
                    pic.crop_bottom = c
                else:             # image too wide -> crop left/right equally
                    keep = target / ar
                    c = (1.0 - keep) / 2.0
                    pic.crop_left = c
                    pic.crop_right = c
    prs.save(out)
    print("wrote %s (%d slides)" % (out, len(imgs)))


if __name__ == "__main__":
    main()
